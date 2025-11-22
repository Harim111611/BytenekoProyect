from django.shortcuts import render, get_object_or_404
from django.http import HttpResponse
from django.utils import timezone
from django.db.models.functions import TruncDate
from django.contrib.auth.decorators import login_required
from surveys.models import Encuesta, RespuestaEncuesta, RespuestaPregunta
from django.core.serializers.json import DjangoJSONEncoder
import json
import collections
import re
import io
import base64
import os
from django.conf import settings
from django.db.models import Count, Avg, Min, Max, StdDev
from datetime import datetime, timedelta
import statistics

from django.template.loader import render_to_string
from django.core.cache import cache

# Imports para PPTX
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Imports Matplotlib
import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns

try:
    from weasyprint import HTML, CSS
except ImportError:
    HTML = None


# ============================================================
# 1. DASHBOARD PRINCIPAL
# ============================================================
@login_required
def dashboard_view(request):
    user = request.user
    cache_key = f"dashboard_data_user_{user.id}"
    context = cache.get(cache_key)

    if not context:
        mis_encuestas = Encuesta.objects.filter(creador=user)
        total_encuestas = mis_encuestas.count()
        activas = mis_encuestas.filter(estado='active').count()
        respuestas_qs = RespuestaEncuesta.objects.filter(encuesta__creador=user)
        total_respuestas = respuestas_qs.count()
        today = timezone.now().date()
        respuestas_hoy = respuestas_qs.filter(creado_en__date=today).count()

        kpis = [
            {
                'label': 'Encuestas Activas',
                'value': activas,
                'icon': 'bi-broadcast',
                'color': 'success',
                'subtext': f'de {total_encuestas} totales'
            },
            {
                'label': 'Respuestas Totales',
                'value': total_respuestas,
                'icon': 'bi-people-fill',
                'color': 'primary',
                'subtext': 'Histórico'
            },
            {
                'label': 'Recibidas Hoy',
                'value': respuestas_hoy,
                'icon': 'bi-bar-chart-fill',
                'color': 'info',
                'subtext': 'Interacciones hoy'
            },
        ]

        start_date = today - timedelta(days=14)
        daily_data = (
            respuestas_qs
            .filter(creado_en__date__gte=start_date)
            .annotate(date=TruncDate('creado_en'))
            .values('date')
            .annotate(c=Count('id'))
            .order_by('date')
        )
        date_map = {item['date']: item['c'] for item in daily_data}
        chart_labels, chart_data = [], []
        curr = start_date
        while curr <= today:
            chart_labels.append(curr.strftime('%d %b'))
            chart_data.append(date_map.get(curr, 0))
            curr += timedelta(days=1)

        status_counts = mis_encuestas.values('estado').annotate(count=Count('id'))
        status_map = {s['estado']: s['count'] for s in status_counts}
        pie_data = [
            status_map.get('active', 0),
            status_map.get('draft', 0),
            status_map.get('closed', 0)
        ]

        recent_surveys_orm = mis_encuestas.annotate(
            response_count=Count('respuestas')
        ).order_by('-fecha_modificacion')[:5]

        recent_activity = []
        for s in recent_surveys_orm:
            resp_count = s.response_count
            goal = s.objetivo_muestra if s.objetivo_muestra > 0 else 1
            progress = int((resp_count / goal) * 100)
            recent_activity.append({
                'id': s.id,
                'titulo': s.titulo,
                'estado': s.estado,
                'get_estado_display': s.get_estado_display(),
                'respuestas': resp_count,
                'objetivo': s.objetivo_muestra,
                'progreso_pct': progress,
                'visual_progress': min(progress, 100),
                'fecha': s.fecha_modificacion
            })

        context = {
            'page_name': 'dashboard',
            'kpis': kpis,
            'chart_labels': json.dumps(chart_labels),
            'chart_data': json.dumps(chart_data),
            'pie_data': json.dumps(pie_data),
            'recent_activity': recent_activity,
            'user_name': request.user.username,
        }
        cache.set(cache_key, context, 300)

    return render(request, 'core/dashboard.html', context)


@login_required
def results_dashboard_view(request):
    user = request.user

    mis_encuestas = Encuesta.objects.filter(creador=user)
    respuestas_qs = RespuestaEncuesta.objects.filter(encuesta__creador=user)

    total_responses = respuestas_qs.count()
    total_active = mis_encuestas.filter(estado='active').count()

    sat_avg = RespuestaPregunta.objects.filter(
        respuesta_encuesta__encuesta__creador=user,
        pregunta__tipo='scale',
        valor_numerico__isnull=False
    ).aggregate(avg=Avg('valor_numerico'))['avg'] or 0

    today = timezone.now().date()
    start_date = today - timedelta(days=30)

    daily_data = (
        respuestas_qs
        .filter(creado_en__date__gte=start_date)
        .annotate(date=TruncDate('creado_en'))
        .values('date')
        .annotate(c=Count('id'))
        .order_by('date')
    )

    date_map = {item['date']: item['c'] for item in daily_data}
    chart_labels = []
    chart_data = []
    curr = start_date
    while curr <= today:
        chart_labels.append(curr.strftime('%d %b'))
        chart_data.append(date_map.get(curr, 0))
        curr += timedelta(days=1)

    top_surveys = mis_encuestas.annotate(
        response_count=Count('respuestas')
    ).order_by('-response_count')[:5]

    context = {
        'page_name': 'results',
        'total_responses': total_responses,
        'total_active': total_active,
        'global_satisfaction': sat_avg,
        'chart_labels': json.dumps(chart_labels),
        'chart_data': json.dumps(chart_data),
        'top_surveys': top_surveys,
    }
    return render(request, 'core/results_dashboard.html', context)


# ============================================================
# 2. HELPERS GRÁFICOS + ANÁLISIS
# ============================================================
def generate_heatmap_chart(df):
    df_numeric = df.select_dtypes(include=['float64', 'int64'])
    if df_numeric.shape[1] < 2 or df_numeric.shape[0] < 2:
        return None

    plt.style.use('default')
    plt.rcParams.update({
        'font.family': 'sans-serif',
        'font.sans-serif': ['Arial'],
        'font.size': 9
    })
    fig, ax = plt.subplots(figsize=(8, 5))
    sns.heatmap(
        df_numeric.corr(),
        annot=True,
        cmap='coolwarm',
        fmt=".2f",
        linewidths=.5,
        vmin=-1,
        vmax=1,
        ax=ax
    )
    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches='tight', transparent=True)
    plt.close(fig)
    buf.seek(0)
    return base64.b64encode(buf.read()).decode("utf-8")


def generate_nps_chart(prom, pas, det):
    labels = ['Promotores', 'Pasivos', 'Detractores']
    sizes = [prom, pas, det]
    colors = ['#10b981', '#fbbf24', '#ef4444']

    labels = [l for l, s in zip(labels, sizes) if s > 0]
    colors = [c for c, s in zip(colors, sizes) if s > 0]
    sizes = [s for s in sizes if s > 0]

    if not sizes:
        return None

    fig, ax = plt.subplots(figsize=(5, 3))
    ax.pie(
        sizes,
        labels=labels,
        autopct='%1.0f%%',
        startangle=90,
        colors=colors,
        pctdistance=0.85,
        textprops=dict(color="#333333", fontsize=10, weight='bold')
    )
    centre_circle = plt.Circle((0, 0), 0.60, fc='white')
    fig.gca().add_artist(centre_circle)
    plt.tight_layout()

    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches='tight', transparent=True)
    plt.close(fig)
    buf.seek(0)
    return base64.b64encode(buf.read()).decode("utf-8")


def generate_vertical_bar_chart(labels, counts, title):
    plt.rcParams.update({'font.family': 'sans-serif', 'font.sans-serif': ['Arial']})
    fig, ax = plt.subplots(figsize=(7, 4))
    bars = ax.bar(labels, counts, color='#0d6efd', alpha=0.9, width=0.6)

    ax.set_title(title, fontsize=12, weight='bold', color='#111827', pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.spines['bottom'].set_color('#e5e7eb')
    ax.tick_params(axis='y', left=False, labelleft=False)
    ax.tick_params(axis='x', colors='#6b7281')
    ax.bar_label(bars, fmt='%d', padding=3, fontsize=10, color='#111827', weight='bold')
    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, transparent=True)
    plt.close(fig)
    buf.seek(0)
    return base64.b64encode(buf.read()).decode("utf-8")


def generate_horizontal_bar_chart(labels, counts, title):
    plt.rcParams.update({'font.family': 'sans-serif', 'font.sans-serif': ['Arial']})
    height = max(3, len(labels) * 0.6)
    fig, ax = plt.subplots(figsize=(7, height))
    y_pos = range(len(labels))
    bars = ax.barh(y_pos, counts, color='#0d6efd', alpha=0.9, height=0.6)

    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, fontsize=10, color='#6b7281')
    ax.invert_yaxis()
    ax.set_title(title, fontsize=12, weight='bold', color='#111827', pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.tick_params(axis='x', bottom=False, labelbottom=False)
    ax.tick_params(axis='y', left=False)
    ax.bar_label(bars, fmt='%d', padding=5, fontsize=10, color='#111827', weight='bold')
    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, transparent=True)
    plt.close(fig)
    buf.seek(0)
    return base64.b64encode(buf.read()).decode("utf-8")


def analizar_texto_avanzado(qs):
    # Extraer textos directamente (qs es queryset de RespuestaPregunta)
    textos = list(qs.values_list('valor_texto', flat=True))
    textos = [t for t in textos if t]

    if not textos: return [], []

    if len(textos) > 2000:
        textos = textos[:2000]

    text_full = " ".join(textos).lower()
    clean = re.sub(r'[^\w\s]', '', text_full)
    words = clean.split()
    stops = {'de', 'la', 'que', 'el', 'en', 'y', 'a', 'los', 'se', 'del', 'las', 'un', 'por', 'con', 'no', 'una', 'su',
             'para', 'es', 'al', 'lo', 'como', 'mas', 'pero', 'sus', 'le', 'ya', 'o', 'muy', 'sin', 'sobre', 'me', 'mi',
             'bueno', 'buena', 'malo', 'mala', 'yo', 'tu'}
    filtered = [w for w in words if w not in stops and len(w) > 2]
    bigrams = [f"{filtered[i]} {filtered[i + 1]}" for i in range(len(filtered) - 1)]
    return collections.Counter(filtered).most_common(5), collections.Counter(bigrams).most_common(3)


def build_responses_df(encuesta, respuestas_qs):
    """
    Versión ULTRA-RÁPIDA para grandes datasets.
    """
    data = RespuestaPregunta.objects.filter(
        respuesta_encuesta__in=respuestas_qs
    ).values(
        'respuesta_encuesta__id',
        'respuesta_encuesta__creado_en',
        'pregunta__texto',
        'valor_texto',
        'valor_numerico',
        'opcion__texto'
    )

    if not data:
        return pd.DataFrame()

    df_raw = pd.DataFrame(list(data))

    if df_raw.empty:
        return pd.DataFrame()

    df_raw['valor'] = df_raw['valor_numerico'].fillna(
        df_raw['opcion__texto']
    ).fillna(df_raw['valor_texto'])

    try:
        df = df_raw.pivot_table(
            index='respuesta_encuesta__id',
            columns='pregunta__texto',
            values='valor',
            aggfunc='first'
        )

        fechas = df_raw.groupby('respuesta_encuesta__id')['respuesta_encuesta__creado_en'].first()
        df['Fecha'] = fechas

        return df
    except Exception as e:
        return pd.DataFrame()


def get_analysis_data_for_report(encuesta, base_respuestas, include_charts=True, cache_key=None):
    if cache_key:
        cached_data = cache.get(cache_key)
        if cached_data:
            return cached_data

    analysis_data = []
    nps_data = {'score': None, 'breakdown_chart': None}

    promedio_satisfaccion = 0
    preguntas_escala = encuesta.preguntas.filter(tipo='scale')
    if preguntas_escala.exists():
        vals = RespuestaPregunta.objects.filter(
            pregunta__in=preguntas_escala,
            respuesta_encuesta__in=base_respuestas,
            valor_numerico__isnull=False
        ).aggregate(avg=Avg('valor_numerico'))
        promedio_satisfaccion = vals['avg'] or 0

    pregunta_nps = preguntas_escala.first()
    if pregunta_nps:
        qs_nps = RespuestaPregunta.objects.filter(
            pregunta=pregunta_nps,
            respuesta_encuesta__in=base_respuestas,
            valor_numerico__isnull=False
        )
        total = qs_nps.count()
        if total > 0:
            prom = qs_nps.filter(valor_numerico__gte=9).count()
            pas = qs_nps.filter(valor_numerico__in=[7, 8]).count()
            det = qs_nps.filter(valor_numerico__lte=6).count()
            nps_data['score'] = round(((prom / total) * 100) - ((det / total) * 100), 1)
            if include_charts:
                nps_data['breakdown_chart'] = generate_nps_chart(prom, pas, det)

    heatmap_image = None
    if include_charts:
        try:
            df = build_responses_df(encuesta, base_respuestas)
            if not df.empty:
                heatmap_image = generate_heatmap_chart(df)
        except Exception:
            pass

    for i, pregunta in enumerate(encuesta.preguntas.all().order_by('orden'), 1):
        item = {
            'id': pregunta.id,
            'orden': i,
            'texto': pregunta.texto,
            'tipo': pregunta.tipo,
            'tipo_display': pregunta.get_tipo_display(),
            'insight': '',
            'chart_image': None,
            'chart_data': None,
            'total_respuestas': 0,
            'estadisticas': None,
            'opciones': [],
            'samples_texto': [],
            'top_options': [],
            'avg': None,
            'scale_cap': None,
        }

        res_preg = RespuestaPregunta.objects.filter(
            pregunta=pregunta,
            respuesta_encuesta__in=base_respuestas
        )
        item['total_respuestas'] = res_preg.count()

        if pregunta.tipo in ['number', 'scale']:
            qs = res_preg.filter(valor_numerico__isnull=False).values_list('valor_numerico', flat=True)
            valores_list = list(qs)

            if valores_list:
                _min = min(valores_list)
                _max = max(valores_list)
                _avg = sum(valores_list) / len(valores_list)
                _med = statistics.median(valores_list)

                item['estadisticas'] = {'minimo': _min, 'maximo': _max, 'promedio': _avg, 'mediana': _med}
                item['avg'] = _avg

                scale_cap = 10 if pregunta.tipo == 'scale' else (_max if _max > 5 else 5)
                item['scale_cap'] = scale_cap
                normalized = (_avg / scale_cap) * 10

                sentimiento = "Neutral"
                if normalized >= 8:
                    sentimiento = "Excelente"
                elif normalized >= 6:
                    sentimiento = "Bueno"
                elif normalized <= 4:
                    sentimiento = "Crítico"
                item[
                    'insight'] = f"<strong>Desempeño {sentimiento}</strong>. Promedio de {_avg:.1f} (Mediana: {_med:.1f})."

                if include_charts:
                    counts = collections.Counter(valores_list)
                    if pregunta.tipo == 'scale':
                        labels = [str(x) for x in range(int(_min), int(_max) + 1)]
                        data = [counts.get(x, 0) for x in range(int(_min), int(_max) + 1)]
                    else:
                        sorted_keys = sorted(counts.keys())
                        labels = [str(k) for k in sorted_keys]
                        data = [counts[k] for k in sorted_keys]

                    item['chart_data'] = {'labels': labels, 'data': data}
                    item['chart_image'] = generate_vertical_bar_chart(labels, data, "Distribución de Respuestas")

        elif pregunta.tipo in ['single', 'multi']:
            all_vals = []
            for r in res_preg:
                if r.opcion:
                    all_vals.append(r.opcion.texto)
                elif r.valor_texto:
                    val_clean = [x.strip() for x in r.valor_texto.split(',') if x.strip()]
                    all_vals.extend(val_clean)

            total_votes = len(all_vals)
            if total_votes > 0:
                counter = collections.Counter(all_vals)
                opciones_list = []
                for label, count in counter.most_common():
                    pct = (count / total_votes) * 100
                    opciones_list.append({'label': label, 'count': count, 'percent': pct})
                item['opciones'] = opciones_list

                top_3 = counter.most_common(3)
                item['top_options'] = top_3
                if top_3:
                    winner, w_count = top_3[0]
                    w_pct = (w_count / total_votes) * 100
                    item[
                        'insight'] = f"Opción líder: <strong>{winner}</strong> con el {w_pct:.0f}% de las preferencias."

                if include_charts:
                    top_chart = counter.most_common(8)
                    labels = [x[0] for x in top_chart]
                    data = [x[1] for x in top_chart]
                    item['chart_data'] = {'labels': labels, 'data': data}
                    item['chart_image'] = generate_horizontal_bar_chart(labels, data, "Top Selecciones")
            else:
                item['opciones'] = []
                item['insight'] = "Aún no hay respuestas para esta pregunta."

        elif pregunta.tipo == 'text':
            qs_text = res_preg.exclude(valor_texto__isnull=True).exclude(valor_texto__exact="")

            qs_text = qs_text.select_related('respuesta_encuesta')

            item['samples_texto'] = list(
                qs_text.order_by('-respuesta_encuesta__creado_en')
                .values_list('valor_texto', flat=True)[:5]
            )

            words, _ = analizar_texto_avanzado(qs_text)
            if words:
                top_kw = ", ".join([w[0] for w in words[:3]])
                item['insight'] = f"Temas recurrentes detectados: <strong>{top_kw}</strong>."
            else:
                item['insight'] = "Respuestas dispersas sin patrones claros aún."

        analysis_data.append(item)

    final_data = {
        'analysis_data': analysis_data,
        'nps_data': nps_data,
        'heatmap_image': heatmap_image,
        'kpi_prom_satisfaccion': promedio_satisfaccion,
    }

    if cache_key:
        cache.set(cache_key, final_data, 3600)

    return final_data


# ============================================================
# 3. VISTAS QUE USAN EL ANÁLISIS
# ============================================================

@login_required
def reports_page_view(request):
    return render(
        request,
        'core/reports_page.html',
        {'page_name': 'reports', 'encuestas': Encuesta.objects.filter(creador=request.user)}
    )


@login_required
def report_preview_ajax(request, pk):
    encuesta = get_object_or_404(Encuesta, pk=pk)
    qs = RespuestaEncuesta.objects.filter(encuesta=encuesta)

    # 1. Leer filtros de frontend
    show_kpis = request.GET.get('include_kpis', 'true') == 'true'
    show_charts = request.GET.get('include_charts', 'true') == 'true'
    show_table = request.GET.get('include_table', 'true') == 'true'

    start = request.GET.get('start_date')
    end = request.GET.get('end_date')
    window = request.GET.get('window_days')

    # 2. Aplicar filtros de fecha
    if start:
        try:
            start_date = datetime.strptime(start, '%Y-%m-%d').date()
            qs = qs.filter(creado_en__date__gte=start_date)
        except ValueError:
            pass
    elif window and window.isdigit():
        days = int(window)
        start_dt = timezone.now() - timedelta(days=days)
        qs = qs.filter(creado_en__gte=start_dt)
        start = start_dt.strftime('%Y-%m-%d')  # Para cache

    if end:
        try:
            end_date = datetime.strptime(end, '%Y-%m-%d').date()
            qs = qs.filter(creado_en__date__lte=end_date)
        except ValueError:
            pass

    cache_key = f"survey_analysis_{encuesta.id}_{start or 'all'}_{end or 'all'}_{window or 'all'}"

    data = get_analysis_data_for_report(encuesta, qs, include_charts=True, cache_key=cache_key)
    json_data = json.dumps(data['analysis_data'], cls=DjangoJSONEncoder)

    return HttpResponse(
        render_to_string(
            'core/_report_preview_content.html',
            {
                'encuesta': encuesta,
                'analysis_data': data['analysis_data'],
                'analysis_data_json': json_data,
                'nps_score': data['nps_data']['score'],
                'total_respuestas': qs.count(),
                'kpi_prom_satisfaccion': data['kpi_prom_satisfaccion'],
                'include_kpis': show_kpis,
                'include_charts': show_charts,
                'include_table': show_table,
            }
        )
    )


@login_required
def report_pdf_view(request):
    survey_id = request.POST.get('survey_id') or request.GET.get('survey_id')
    if not survey_id:
        return HttpResponse("ID de encuesta faltante", status=400)

    encuesta = get_object_or_404(Encuesta, pk=survey_id)
    qs = RespuestaEncuesta.objects.filter(encuesta=encuesta)

    start = request.POST.get('start_date') or request.GET.get('start_date')
    end = request.POST.get('end_date') or request.GET.get('end_date')
    window = request.POST.get('window_days')
    include_table = request.POST.get('include_table') == 'on'

    if start:
        start_date = datetime.strptime(start, '%Y-%m-%d').date()
        qs = qs.filter(creado_en__date__gte=start_date)
    elif window and window.isdigit():
        days = int(window)
        start_dt = timezone.now() - timedelta(days=days)
        qs = qs.filter(creado_en__gte=start_dt)
        start = start_dt.strftime('%Y-%m-%d')

    if end:
        end_date = datetime.strptime(end, '%Y-%m-%d').date()
        qs = qs.filter(creado_en__date__lte=end_date)

    cache_key = f"survey_analysis_{encuesta.id}_{start or 'all'}_{end or 'all'}_{window or 'all'}"

    data_pack = get_analysis_data_for_report(encuesta, qs, include_charts=True, cache_key=cache_key)

    context = {
        'encuesta': encuesta,
        'start_date': start,
        'end_date': end,
        'total_respuestas': qs.count(),
        'analysis_data': data_pack['analysis_data'],
        'nps_score': data_pack['nps_data']['score'],
        'kpi_prom_satisfaccion': data_pack['kpi_prom_satisfaccion'],
        'include_table': include_table,
    }

    html_string = render_to_string('core/report_pdf_template.html', context)

    if not HTML:
        return HttpResponse("Error: Librería WeasyPrint no instalada.", status=500)

    try:
        pdf_file = HTML(string=html_string, base_url=request.build_absolute_uri()).write_pdf()
    except Exception as e:
        return HttpResponse(f"Error al generar PDF: {e}", status=500)

    response = HttpResponse(pdf_file, content_type='application/pdf')
    filename = f"Reporte_{encuesta.titulo[:20]}_{datetime.now().strftime('%Y%m%d')}.pdf"
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response


@login_required
def report_powerpoint_view(request):
    if request.method != 'POST':
        return HttpResponse("Método no permitido.", status=405)

    survey_id = request.POST.get('survey_id')
    encuesta = get_object_or_404(Encuesta, pk=survey_id)

    qs_all = RespuestaEncuesta.objects.filter(encuesta=encuesta)
    qs = qs_all

    start = request.POST.get('start_date')
    end = request.POST.get('end_date')
    window = request.POST.get('window_days')
    date_range_label = "Todo el histórico"

    if start:
        start_date = datetime.strptime(start, '%Y-%m-%d').date()
        qs = qs.filter(creado_en__date__gte=start_date)
        date_range_label = f"Desde {start_date.strftime('%d/%m/%Y')}"
    elif window and window.isdigit():
        days = int(window)
        start_dt = timezone.now() - timedelta(days=days)
        qs = qs.filter(creado_en__gte=start_dt)
        date_range_label = f"Últimos {days} días"
        start = start_dt.strftime('%Y-%m-%d')

    if end:
        end_date = datetime.strptime(end, '%Y-%m-%d').date()
        qs = qs.filter(creado_en__date__lte=end_date)
        date_range_label += f" hasta {end_date.strftime('%d/%m/%Y')}"

    cache_key = f"survey_analysis_{encuesta.id}_{start or 'all'}_{end or 'all'}_{window or 'all'}"

    data = get_analysis_data_for_report(encuesta, qs, include_charts=True, cache_key=cache_key)
    analysis_data = data['analysis_data']
    nps_data = data['nps_data']
    nps_score = nps_data['score']
    heatmap_image = data['heatmap_image']

    numeric_scores = []
    sat_sum_10 = 0.0
    sat_q_count = 0
    for item in analysis_data:
        avg_val = item.get('avg')
        scale_cap = item.get('scale_cap')
        if avg_val is not None and scale_cap:
            norm_10 = (avg_val / scale_cap) * 10
            numeric_scores.append({
                'orden': item['orden'],
                'texto': item['texto'],
                'norm_10': norm_10,
            })
            sat_sum_10 += norm_10
            sat_q_count += 1

    avg_sat = round(sat_sum_10 / sat_q_count, 1) if sat_q_count else None

    strengths = []
    opportunities = []
    if numeric_scores:
        numeric_desc = sorted(numeric_scores, key=lambda x: x['norm_10'], reverse=True)
        numeric_asc = sorted(numeric_scores, key=lambda x: x['norm_10'])

        strengths = [n for n in numeric_desc if n['norm_10'] >= 8][:3]
        opportunities = [n for n in numeric_asc if n['norm_10'] < 7][:3]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BYTE_BLUE = RGBColor(13, 110, 253)
    BYTE_TEXT = RGBColor(17, 24, 39)
    BYTE_GRAY = RGBColor(107, 114, 129)
    BYTE_BG_CARD = RGBColor(248, 249, 250)
    BYTE_BORDER = RGBColor(222, 226, 230)

    logo_path = os.path.join(settings.BASE_DIR, 'static', 'img', 'byteneko_logo.png')
    has_logo = os.path.exists(logo_path)

    def apply_header(slide, title_text):
        bar_height = Inches(1.0)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, bar_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = BYTE_BLUE
        bar.line.fill.background()

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.96), prs.slide_width, Inches(0.04))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(15, 118, 110)
        accent.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0), Inches(10), bar_height)
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(24)
        p.font.bold = True

        if has_logo:
            try:
                pic_height = Inches(0.6)
                top_pos = (bar_height - pic_height) / 2
                slide.shapes.add_picture(logo_path, prs.slide_width - Inches(2.0), top_pos, height=pic_height)
            except:
                pass
        else:
            tb_logo = slide.shapes.add_textbox(prs.slide_width - Inches(2.5), 0, Inches(2.0), bar_height)
            tf_logo = tb_logo.text_frame
            tf_logo.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_logo = tf_logo.paragraphs[0]
            p_logo.text = "BYTENEKO"
            p_logo.font.color.rgb = RGBColor(255, 255, 255)
            p_logo.font.bold = True
            p_logo.alignment = PP_ALIGN.RIGHT

    def draw_kpi_card(slide, x, y, title, value, color=BYTE_BLUE, subtitle=None):
        value = "--" if value is None else str(value)
        w_card = Inches(3.6)
        h_card = Inches(1.9)
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.05), y + Inches(0.05), w_card, h_card)
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = RGBColor(210, 210, 210)
        shadow.line.fill.background()

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w_card, h_card)
        card.fill.solid()
        card.fill.fore_color.rgb = BYTE_BG_CARD
        card.line.color.rgb = BYTE_BORDER

        tb_title = slide.shapes.add_textbox(x, y + Inches(0.15), w_card, Inches(0.4))
        tf_title = tb_title.text_frame
        tf_title.vertical_anchor = MSO_ANCHOR.TOP
        p1 = tf_title.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(13)
        p1.font.color.rgb = BYTE_GRAY
        p1.alignment = PP_ALIGN.CENTER

        tb_val = slide.shapes.add_textbox(x, y + Inches(0.4), w_card, Inches(1.0))
        tf_val = tb_val.text_frame
        tf_val.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf_val.paragraphs[0]
        p2.text = value
        p2.font.size = Pt(36)
        p2.font.bold = True
        p2.font.color.rgb = color
        p2.alignment = PP_ALIGN.CENTER

        if subtitle:
            tb_sub = slide.shapes.add_textbox(x, y + Inches(1.4), w_card, Inches(0.4))
            tf_sub = tb_sub.text_frame
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle
            p_sub.font.size = Pt(10)
            p_sub.font.color.rgb = BYTE_GRAY
            p_sub.alignment = PP_ALIGN.CENTER

    def add_footer(slide, page_number, total_pages):
        tb = slide.shapes.add_textbox(Inches(0.6), prs.slide_height - Inches(0.45), prs.slide_width - Inches(1.2),
                                      Inches(0.3))
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        p = tf.paragraphs[0]
        p.text = f"Página {page_number} de {total_pages} • {date_range_label} • Byteneko Analytics v1.0"
        p.font.size = Pt(9)
        p.font.color.rgb = BYTE_GRAY
        p.alignment = PP_ALIGN.LEFT

    def truncate_title(text, max_len=60):
        text = text or ""
        return text if len(text) <= max_len else text[: max_len - 3] + "..."

    def split_question_title(texto):
        texto = texto or ""
        if "(" in texto and ")" in texto:
            base, extra = texto.split("(", 1)
            extra = "(" + extra
        else:
            base, extra = texto, ""
        return base.strip(), extra.strip()

    def is_text_like_question(item):
        texto = (item.get('texto') or '').lower()
        if item.get('tipo') == 'text': return True
        keywords = ['comentario', 'sugerencia', 'observación']
        return any(k in texto for k in keywords)

    total_pages = 4 + len(analysis_data) + (1 if heatmap_image else 0)
    current_page = 1

    # SLIDE 1: PORTADA
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    side_width = Inches(4.0)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, side_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BYTE_BLUE
    bg.line.fill.background()

    if has_logo:
        try:
            logo_w = Inches(2.5)
            left_pos = (side_width - logo_w) / 2
            slide.shapes.add_picture(logo_path, left_pos, Inches(1.0), width=logo_w)
        except:
            pass

    right_area_w = prs.slide_width - side_width
    right_area_start = side_width
    textbox_h = Inches(4.0)
    textbox_y = (prs.slide_height - textbox_h) / 2

    tb_title = slide.shapes.add_textbox(right_area_start + Inches(0.5), textbox_y, right_area_w - Inches(1.0),
                                        textbox_h)
    tf_title = tb_title.text_frame
    tf_title.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_title.word_wrap = True

    titulo_encuesta = encuesta.titulo or ""
    if titulo_encuesta.lower().endswith('.csv'):
        titulo_encuesta = os.path.splitext(titulo_encuesta)[0]

    p = tf_title.paragraphs[0]
    p.text = titulo_encuesta
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BYTE_TEXT
    p.alignment = PP_ALIGN.LEFT

    p_sub = tf_title.add_paragraph()
    p_sub.text = f"\nReporte generado: {datetime.now().strftime('%d/%m/%Y')}"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = BYTE_GRAY

    p_sub2 = tf_title.add_paragraph()
    p_sub2.text = f"Periodo: {date_range_label}"
    p_sub2.font.size = Pt(16)
    p_sub2.font.color.rgb = BYTE_GRAY

    tb_tag = slide.shapes.add_textbox(right_area_start + Inches(0.5), prs.slide_height - Inches(0.8),
                                      right_area_w - Inches(1.0), Inches(0.5))
    p_tag = tb_tag.text_frame.paragraphs[0]
    p_tag.text = "Byteneko Analytics · Reporte automático"
    p_tag.font.size = Pt(11)
    p_tag.font.color.rgb = BYTE_GRAY
    p_tag.alignment = PP_ALIGN.RIGHT

    first_response_dt = qs_all.aggregate(first=Min('creado_en'))['first']
    description_lines = []
    raw_desc = (getattr(encuesta, "descripcion", "") or "").strip()
    if raw_desc:
        description_lines.append(raw_desc)
    elif first_response_dt:
        description_lines.append(f"Importada {first_response_dt.date().isoformat()}")
    if not description_lines:
        description_lines.append("Sin descripción disponible.")

    add_footer(slide, current_page, total_pages)
    current_page += 1

    # SLIDE 2: AGENDA
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_header(slide, "Agenda del reporte")
    content_margin = Inches(1.5)
    content_width = prs.slide_width - (content_margin * 2)
    tb_agenda = slide.shapes.add_textbox(content_margin, Inches(2.0), content_width, Inches(4.0))
    tf_ag = tb_agenda.text_frame
    tf_ag.word_wrap = True

    p0 = tf_ag.paragraphs[0]
    p0.text = "En este reporte encontrarás:"
    p0.font.size = Pt(20)
    p0.font.bold = True
    p0.font.color.rgb = BYTE_TEXT
    p0.space_after = Pt(20)

    bullets = [
        "Resumen ejecutivo con KPIs clave.",
        "Mapa de relaciones entre variables (heatmap).",
        "Conclusiones clave y principales oportunidades de mejora.",
        "Detalle por pregunta con gráficas y hallazgos específicos."
    ]
    for txt in bullets:
        pb = tf_ag.add_paragraph()
        pb.text = f"•  {txt}"
        pb.font.size = Pt(16)
        pb.font.color.rgb = BYTE_GRAY
        pb.space_after = Pt(14)

    add_footer(slide, current_page, total_pages)
    current_page += 1

    # SLIDE 3: RESUMEN
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_header(slide, "Resumen ejecutivo")

    card_w = Inches(3.5)
    card_h = Inches(1.85)
    gap = Inches(0.5)
    total_group_w = (card_w * 3) + (gap * 2)
    start_x = (prs.slide_width - total_group_w) / 2

    y_desc = Inches(1.3)
    h_desc = Inches(0.9)
    y_kpi = Inches(2.5)
    y_chart_label = Inches(4.6)
    y_chart_img = Inches(4.9)
    h_chart_target = Inches(2.2)

    card_desc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x, y_desc, total_group_w, h_desc)
    card_desc.fill.solid()
    card_desc.fill.fore_color.rgb = BYTE_BG_CARD
    card_desc.line.color.rgb = BYTE_BORDER

    tb_desc = slide.shapes.add_textbox(start_x + Inches(0.2), y_desc, total_group_w - Inches(0.4), h_desc)
    tf = tb_desc.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_label = tf.paragraphs[0]
    p_label.text = "Sobre esta encuesta: "
    p_label.font.bold = True
    p_label.font.size = Pt(12)
    p_label.font.color.rgb = BYTE_BLUE
    run = p_label.add_run()
    run.text = " ".join(description_lines)
    run.font.bold = False
    run.font.color.rgb = BYTE_TEXT

    draw_kpi_card(slide, start_x, y_kpi, "Total respuestas", qs.count())

    nps_text = None if nps_score is None else f"{nps_score:.1f}"
    nps_color = BYTE_BLUE
    if nps_score is not None:
        if nps_score >= 50:
            nps_color = RGBColor(16, 185, 129)
        elif nps_score < 0:
            nps_color = RGBColor(239, 68, 68)
    draw_kpi_card(slide, start_x + card_w + gap, y_kpi, "NPS global", nps_text, nps_color)

    avg_sat_text = None if avg_sat is None else f"{avg_sat:.1f}"
    draw_kpi_card(slide, start_x + (card_w + gap) * 2, y_kpi, "Satisfacción (/10)", avg_sat_text)

    if nps_data.get('breakdown_chart'):
        tb_g = slide.shapes.add_textbox(start_x, y_chart_label, total_group_w, Inches(0.3))
        p_g = tb_g.text_frame.paragraphs[0]
        p_g.text = "Distribución de sentimiento"
        p_g.alignment = PP_ALIGN.CENTER
        p_g.font.size = Pt(11)
        p_g.font.bold = True
        p_g.font.color.rgb = BYTE_GRAY

        chart_img = io.BytesIO(base64.b64decode(nps_data['breakdown_chart']))
        aspect_ratio = 5.0 / 3.0
        chart_width = h_chart_target * aspect_ratio
        chart_x = (prs.slide_width - chart_width) / 2
        slide.shapes.add_picture(chart_img, chart_x, y_chart_img, width=chart_width, height=h_chart_target)

    add_footer(slide, current_page, total_pages)
    current_page += 1

    # SLIDE HEATMAP
    if heatmap_image:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_header(slide, "Mapa de calor (Correlaciones)")
        img = io.BytesIO(base64.b64decode(heatmap_image))
        margin = Inches(0.8)
        available_w = prs.slide_width - (margin * 2)
        img_w = available_w * 0.6
        img_x = margin
        img_y = Inches(1.8)
        slide.shapes.add_picture(img, img_x, img_y, width=img_w)

        text_x = img_x + img_w + Inches(0.5)
        text_w = available_w * 0.35
        tb = slide.shapes.add_textbox(text_x, Inches(2.5), text_w, Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "¿Cómo leer esto?"
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = BYTE_BLUE
        p2 = tf.add_paragraph()
        p2.text = (
            "Este gráfico muestra qué preguntas están relacionadas.\n\n"
            "• Rojo (cerca de 1): Si una sube, la otra también.\n"
            "• Azul (cerca de -1): Relación inversa.\n"
        )
        p2.font.size = Pt(12)
        p2.font.color.rgb = BYTE_TEXT
        add_footer(slide, current_page, total_pages)
        current_page += 1

    # SLIDE CONCLUSIONES
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_header(slide, "Conclusiones y oportunidades")
    col_margin = Inches(0.8)
    col_gap = Inches(0.6)
    col_w = (prs.slide_width - (col_margin * 2) - col_gap) / 2
    col_y = Inches(1.5)
    col_h = Inches(4.8)

    bg1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_margin, col_y, col_w, col_h)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = BYTE_BG_CARD
    bg1.line.color.rgb = BYTE_BLUE
    bg1.line.width = Pt(2)

    tb1 = slide.shapes.add_textbox(col_margin + Inches(0.2), col_y + Inches(0.2), col_w - Inches(0.4),
                                   col_h - Inches(0.4))
    tf1 = tb1.text_frame
    p_t1 = tf1.paragraphs[0]
    p_t1.text = "Top Fortalezas"
    p_t1.font.bold = True
    p_t1.font.size = Pt(18)
    p_t1.font.color.rgb = BYTE_BLUE
    p_t1.alignment = PP_ALIGN.CENTER
    p_t1.space_after = Pt(15)

    if strengths:
        for s in strengths:
            p = tf1.add_paragraph()
            p.text = f"• P{s['orden']}: {truncate_title(s['texto'], 45)}"
            p.font.bold = True
            p.font.size = Pt(12)
            p.space_before = Pt(10)
            p_score = tf1.add_paragraph()
            p_score.text = f"   Puntaje: {s['norm_10']:.1f}/10"
            p_score.font.size = Pt(12)
            p_score.font.color.rgb = BYTE_GRAY
    else:
        p = tf1.add_paragraph()
        p.text = "Faltan datos para identificar fortalezas."
        p.alignment = PP_ALIGN.CENTER

    x_col2 = col_margin + col_w + col_gap
    bg2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_col2, col_y, col_w, col_h)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = BYTE_BG_CARD
    bg2.line.color.rgb = RGBColor(220, 38, 38)
    bg2.line.width = Pt(2)

    tb2 = slide.shapes.add_textbox(x_col2 + Inches(0.2), col_y + Inches(0.2), col_w - Inches(0.4), col_h - Inches(0.4))
    tf2 = tb2.text_frame
    p_t2 = tf2.paragraphs[0]
    p_t2.text = "Áreas de Mejora"
    p_t2.font.bold = True
    p_t2.font.size = Pt(18)
    p_t2.font.color.rgb = RGBColor(220, 38, 38)
    p_t2.alignment = PP_ALIGN.CENTER
    p_t2.space_after = Pt(15)

    if opportunities:
        for o in opportunities:
            p = tf2.add_paragraph()
            p.text = f"• P{o['orden']}: {truncate_title(o['texto'], 45)}"
            p.font.bold = True
            p.font.size = Pt(12)
            p.space_before = Pt(10)
            p_score = tf2.add_paragraph()
            p_score.text = f"   Puntaje: {o['norm_10']:.1f}/10"
            p_score.font.size = Pt(12)
            p_score.font.color.rgb = BYTE_GRAY
    else:
        p = tf2.add_paragraph()
        p.text = "No se detectan áreas críticas."
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide, current_page, total_pages)
    current_page += 1

    # SLIDES DETALLE
    for item in analysis_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        base_title, extra = split_question_title(item['texto'])
        header_title = f"P{item['orden']}. {truncate_title(base_title)}"
        apply_header(slide, header_title)

        if extra or len(base_title) > 50:
            full_text = (base_title + " " + extra).strip()
            tb_full = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12), Inches(0.5))
            p_full = tb_full.text_frame.paragraphs[0]
            p_full.text = full_text
            p_full.font.color.rgb = BYTE_GRAY
            p_full.font.italic = True
            p_full.font.size = Pt(11)

        text_like = is_text_like_question(item)
        content_start_y = Inches(1.6)

        if item['chart_image'] and not text_like:
            margin = Inches(0.6)
            gap = Inches(0.4)
            w_chart_area = Inches(7.5)
            h_chart_area = Inches(4.8)

            bg_chart = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, margin, content_start_y, w_chart_area,
                                              h_chart_area)
            bg_chart.fill.solid()
            bg_chart.fill.fore_color.rgb = RGBColor(252, 252, 253)
            bg_chart.line.color.rgb = BYTE_BORDER

            img = io.BytesIO(base64.b64decode(item['chart_image']))
            slide.shapes.add_picture(img, margin + Inches(0.2), content_start_y + Inches(0.2),
                                     width=w_chart_area - Inches(0.4))

            x_panel = margin + w_chart_area + gap
            w_panel = prs.slide_width - x_panel - margin

            card_h = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_panel, content_start_y, w_panel,
                                            h_chart_area)
            card_h.fill.solid()
            card_h.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card_h.line.color.rgb = BYTE_BLUE

            tb_h = slide.shapes.add_textbox(x_panel + Inches(0.2), content_start_y + Inches(0.2), w_panel - Inches(0.4),
                                            h_chart_area - Inches(0.4))
            tf_h = tb_h.text_frame
            tf_h.word_wrap = True

            p_t = tf_h.paragraphs[0]
            p_t.text = "Análisis"
            p_t.font.bold = True
            p_t.font.color.rgb = BYTE_BLUE
            p_t.font.size = Pt(14)
            p_t.space_after = Pt(10)

            clean_insight = re.sub(r'<[^>]+>', '', item['insight']) or "Sin datos relevantes."
            p_ins = tf_h.add_paragraph()
            p_ins.text = clean_insight
            p_ins.font.size = Pt(12)
            p_ins.space_after = Pt(20)

            if item.get('top_options'):
                p_top = tf_h.add_paragraph()
                p_top.text = "Top Respuestas:"
                p_top.font.bold = True
                p_top.font.size = Pt(12)
                for l, v in item['top_options']:
                    p_opt = tf_h.add_paragraph()
                    p_opt.text = f"• {l} ({v})"
                    p_opt.font.size = Pt(11)

        else:
            margin_text = Inches(1.5)
            w_text_area = prs.slide_width - (margin_text * 2)
            tb_t = slide.shapes.add_textbox(margin_text, content_start_y, w_text_area, Inches(4.5))
            tf_t = tb_t.text_frame

            p_head = tf_t.paragraphs[0]
            p_head.text = "Resumen de respuestas abiertas"
            p_head.font.bold = True
            p_head.font.size = Pt(16)
            p_head.font.color.rgb = BYTE_BLUE
            p_head.alignment = PP_ALIGN.CENTER
            p_head.space_after = Pt(20)

            insight_text = re.sub(r'<[^>]+>', '', item['insight'])
            p_body = tf_t.add_paragraph()
            p_body.text = insight_text
            p_body.alignment = PP_ALIGN.CENTER
            p_body.font.size = Pt(14)

            comentarios = RespuestaPregunta.objects.filter(
                pregunta_id=item['id'], respuesta_encuesta__in=qs, valor_texto__isnull=False
            ).exclude(valor_texto__exact="")

            # Optimización: select_related para evitar N+1
            comentarios = comentarios.select_related('respuesta_encuesta').order_by('-respuesta_encuesta__creado_en')[
                :3]

            if comentarios:
                p_ex = tf_t.add_paragraph()
                p_ex.text = "\nAlgunos comentarios literales:"
                p_ex.font.bold = True
                p_ex.font.color.rgb = BYTE_GRAY
                p_ex.space_before = Pt(20)

                for c in comentarios:
                    p_c = tf_t.add_paragraph()
                    p_c.text = f"“{(c.valor_texto or '').strip()}”"
                    p_c.font.italic = True
                    p_c.font.size = Pt(12)
                    p_c.space_before = Pt(5)

        add_footer(slide, current_page, total_pages)
        current_page += 1

    f = io.BytesIO()
    prs.save(f)
    f.seek(0)
    response = HttpResponse(
        f.getvalue(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    response['Content-Disposition'] = f'attachment; filename="Reporte_{encuesta.id}.pptx"'
    return response


@login_required
def settings_view(request):
    return render(
        request,
        'core/settings.html',
        {'page_name': 'settings', 'user': request.user}
    )