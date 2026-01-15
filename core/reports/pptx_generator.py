"""core/reports/pptx_generator.py"""
import base64
import io
import logging
from dataclasses import dataclass
from typing import Tuple, Optional, Any, cast, List, Dict

from pptx import Presentation as make_presentation
from pptx.presentation import Presentation as PresentationType
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from django.conf import settings
from django.utils import timezone

logger = logging.getLogger(__name__)


def _apply_text_fitting(
    text_frame,
    *,
    vertical_anchor: Any = MSO_ANCHOR.TOP,
    margin: Any = Inches(0.06),
    auto_size: Any = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
) -> None:
    """Aplica ajustes para evitar que el texto se salga del contenedor.

    Nota: `python-pptx` no recalcula siempre el layout al guardar, pero al abrir
    el archivo en PowerPoint/LibreOffice normalmente respeta `auto_size`.
    """
    try:
        text_frame.word_wrap = True
    except Exception:
        pass
    try:
        text_frame.auto_size = auto_size
    except Exception:
        pass
    try:
        text_frame.vertical_anchor = vertical_anchor
    except Exception:
        pass
    # Reducir márgenes internos ayuda a que el texto quepa mejor.
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        try:
            setattr(text_frame, attr, margin)
        except Exception:
            pass


def _soft_wrap_text(text: str, *, max_run: int = 26) -> str:
    """Inserta oportunidades de corte de línea para tokens largos sin espacios.

    PowerPoint suele respetar \u200b (zero-width space) como punto de quiebre.
    """
    if not text:
        return text
    out: List[str] = []
    run = 0
    for ch in str(text):
        out.append(ch)
        if ch.isspace() or ch in "\n\r\t":
            run = 0
            continue
        run += 1
        if ch in "_-/.:,;" or run >= max_run:
            out.append("\u200b")
            run = 0
    return "".join(out)


def _is_text_like_question(item: dict) -> bool:
    q_type = (item.get('type') or '').lower()
    if q_type in {'text', 'comment', 'comments', 'textarea', 'open'}:
        return True
    if (item.get('tipo_display') or '').lower() == 'text':
        return True
    insight = item.get('insight_data') or {}
    if isinstance(insight, dict) and (insight.get('type') or '').lower() == 'text':
        return True
    return False


def _add_picture_contain(slide, img_bytes: bytes, x, y, w, h, *, padding: Any = Inches(0.1)) -> None:
    """Inserta una imagen dentro de un rectángulo sin deformarla.

    Usa un ajuste tipo "contain": mantiene proporción, centra y respeta padding.
    """
    inner_x = x + padding
    inner_y = y + padding
    inner_w = w - (padding * 2)
    inner_h = h - (padding * 2)

    if int(inner_w) <= 0 or int(inner_h) <= 0:
        return

    try:
        from PIL import Image

        with Image.open(io.BytesIO(img_bytes)) as im:
            px_w, px_h = im.size
        if not px_w or not px_h:
            return

        box_w = float(int(inner_w))
        box_h = float(int(inner_h))
        img_ar = float(px_w) / float(px_h)
        box_ar = box_w / box_h

        if img_ar >= box_ar:
            target_w = int(inner_w)
            target_h = int(round(target_w / img_ar))
        else:
            target_h = int(inner_h)
            target_w = int(round(target_h * img_ar))

        left = int(inner_x) + int(round((int(inner_w) - target_w) / 2))
        top = int(inner_y) + int(round((int(inner_h) - target_h) / 2))

        slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=target_w, height=target_h)
        return
    except Exception:
        # Fallback: no deformar; preferir ajustar por ancho.
        try:
            pic = slide.shapes.add_picture(io.BytesIO(img_bytes), inner_x, inner_y, width=inner_w)
            if int(pic.height) > int(inner_h):
                # Si queda muy alta, reintentar ajustando por alto.
                slide.shapes._spTree.remove(pic._element)
                slide.shapes.add_picture(io.BytesIO(img_bytes), inner_x, inner_y, height=inner_h)
        except Exception:
            logger.exception("Error insertando imagen en PPTX")


class PPTXTheme:
    """Definición centralizada de estilos y colores para el reporte."""
    
    # Colores Principales
    BRAND_BLUE = RGBColor(0, 80, 158)       # Azul Byteneko principal
    ACCENT_COLOR = RGBColor(79, 70, 229)    # Acento (#4F46E5)
    DARK_INK = RGBColor(17, 24, 39)         # Texto principal (#111827)
    MUTED_INK = RGBColor(107, 114, 128)     # Texto secundario (#6B7280)
    
    # Colores de Fondo y UI
    BG_LIGHT = RGBColor(249, 250, 251)      # Fondo general muy suave
    WHITE = RGBColor(255, 255, 255)
    BORDER_LIGHT = RGBColor(229, 231, 235)
    
    # Colores Semánticos
    SUCCESS = RGBColor(16, 185, 129)        # Verde positivo
    DANGER = RGBColor(239, 68, 68)          # Rojo negativo
    WARNING = RGBColor(245, 158, 11)        # Naranja alerta
    
    # Tipografía
    FONT_FAMILY = "Segoe UI"
    FONT_TITLE_SIZE = Pt(24) # Reduced from 28 to prevent overflow
    FONT_BODY_SIZE = Pt(11)  # Reduced from 12
    
    # Dimensiones Layout
    MARGIN_X = Inches(0.5)
    MARGIN_Y = Inches(0.5)
    SLIDE_WIDTH = Inches(13.33)
    SLIDE_HEIGHT = Inches(7.5)

    @staticmethod
    def mix_color(rgb: RGBColor, opacity: float) -> RGBColor:
        """Simula transparencia mezclando con blanco."""
        w = max(0.0, min(1.0, float(opacity)))
        r = int(round(255 * (1 - w) + rgb[0] * w))
        g = int(round(255 * (1 - w) + rgb[1] * w))
        b = int(round(255 * (1 - w) + rgb[2] * w))
        return RGBColor(r, g, b)


def _safe_slide_dims(prs: PresentationType) -> Tuple[Any, Any]:
    w = prs.slide_width if prs.slide_width else PPTXTheme.SLIDE_WIDTH
    h = prs.slide_height if prs.slide_height else PPTXTheme.SLIDE_HEIGHT
    return cast(Any, w), cast(Any, h)


class SlideComposer:
    """Helper para componer elementos en una diapositiva con estilo consistente."""
    
    def __init__(self, slide, width, height):
        self.slide = slide
        self.width = width
        self.height = height
        self.theme = PPTXTheme

    def add_background_shape(self):
        """Añade un fondo sutil a la diapositiva."""
        bg = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.width, self.height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.theme.BG_LIGHT
        bg.line.fill.background()
        self._send_to_back(bg)

    def add_header(self, title: str, subtitle: str = ""):
        """Crea un encabezado moderno y limpio."""
        # Banda superior sutil o barra de acento
        bar = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), self.height
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme.BRAND_BLUE
        bar.line.fill.background()

        # Título
        left = Inches(0.5)
        top = Inches(0.4)
        width = self.width - Inches(1.0)
        height = Inches(0.8)

        tb = self.slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        _apply_text_fitting(tf)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.theme.FONT_FAMILY
        p.font.size = self.theme.FONT_TITLE_SIZE
        p.font.bold = True
        p.font.color.rgb = self.theme.DARK_INK

        
        if subtitle:
            p_sub = tf.add_paragraph()
            p_sub.text = subtitle
            p_sub.font.name = self.theme.FONT_FAMILY
            p_sub.font.size = Pt(14)
            p_sub.font.color.rgb = self.theme.MUTED_INK

    def add_card(self, left, top, width, height, bg_color=None) -> Any:
        """Añade una 'tarjeta' contenedora con sombra suave simulada."""
        # Sombra offset
        shadow = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left + Inches(0.04), top + Inches(0.04), width, height
        )
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = RGBColor(220, 220, 225)
        shadow.line.fill.background()
        
        # Tarjeta real
        card = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color or self.theme.WHITE
        card.line.color.rgb = self.theme.BORDER_LIGHT
        card.line.width = Pt(0.75)
        
        return card

    def add_metric_big(self, left, top, label: str, value: str, trend_delta: Optional[float] = None, width: Any = None):
        """Añade un bloque de métrica grande destacado."""
        container_width = width or Inches(3.0)
        
        # Label
        tb_lbl = self.slide.shapes.add_textbox(left, top, container_width, Inches(0.3))
        _apply_text_fitting(tb_lbl.text_frame, margin=Inches(0.03))
        p_lbl = tb_lbl.text_frame.paragraphs[0]
        p_lbl.text = label.upper()
        p_lbl.font.size = Pt(10)
        p_lbl.font.color.rgb = self.theme.MUTED_INK
        p_lbl.font.bold = True
        p_lbl.font.name = self.theme.FONT_FAMILY

        # Value - Ajuste dinámico de tamaño de fuente con AUTOSIZE
        val_height = Inches(0.8)
        tb_val = self.slide.shapes.add_textbox(left, top + Inches(0.25), container_width, val_height)
        tf_val = tb_val.text_frame
        _apply_text_fitting(tf_val)
        
        p_val = tf_val.paragraphs[0]
        p_val.text = value
        p_val.font.size = Pt(40) # Tamaño base ideal
        p_val.font.bold = True
        p_val.font.color.rgb = self.theme.BRAND_BLUE
        p_val.font.name = self.theme.FONT_FAMILY
        
        # Trend
        if trend_delta is not None and abs(trend_delta) > 0.1:
            text_len = len(str(value))
            tb_trend = self.slide.shapes.add_textbox(left, top + Inches(1.15 if text_len > 15 else 0.95), container_width, Inches(0.3))
            _apply_text_fitting(tb_trend.text_frame, margin=Inches(0.03))
            p_trend = tb_trend.text_frame.paragraphs[0]
            symbol = "▲" if trend_delta > 0 else "▼"
            color = self.theme.SUCCESS if trend_delta > 0 else self.theme.DANGER
            p_trend.text = f"{symbol} {abs(trend_delta):.1f}% vs. periodo ant."
            p_trend.font.size = Pt(11)
            p_trend.font.color.rgb = color
            p_trend.font.name = self.theme.FONT_FAMILY

    def add_narrative_box(self, left, top, width, height, text: str):
        """Caja para narrativa/insight IA con estilo distintivo."""
        bg_color = self.theme.mix_color(self.theme.BRAND_BLUE, 0.05)
        border_color = self.theme.mix_color(self.theme.BRAND_BLUE, 0.2)
        
        box = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.color.rgb = border_color
        box.line.width = Pt(1)
        
        # Padding interno simulado
        inset_x = Inches(0.2)
        inset_y = Inches(0.15)
        
        # Icono/Label "IA Insight"
        tb_label = self.slide.shapes.add_textbox(left + inset_x, top + inset_y, width - (inset_x*2), Inches(0.3))
        _apply_text_fitting(tb_label.text_frame, margin=Inches(0.03))
        p_lbl = tb_label.text_frame.paragraphs[0]
        p_lbl.text = "ANALISIS INTELIGENTE"
        p_lbl.font.size = Pt(9)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = self.theme.BRAND_BLUE
        p_lbl.font.name = self.theme.FONT_FAMILY
        
        # Contenido
        content_top = top + inset_y + Inches(0.25)
        tb_text = self.slide.shapes.add_textbox(left + inset_x, content_top, width - (inset_x*2), height - (inset_y*2) - Inches(0.25))
        tf = tb_text.text_frame
        _apply_text_fitting(tf)
        
        p_txt = tf.paragraphs[0]
        p_txt.text = text
        p_txt.font.size = Pt(12) # Base size, will shrink
        p_txt.font.color.rgb = self.theme.DARK_INK
        p_txt.font.italic = True
        p_txt.font.name = self.theme.FONT_FAMILY
        p_txt.alignment = PP_ALIGN.JUSTIFY

    def _send_to_back(self, shape):
        """Mueve una forma al fondo del z-order."""
        try:
            sp_tree = self.slide.shapes._spTree
            sp_tree.remove(shape._element)
            sp_tree.insert(2, shape._element)
        except Exception:
            pass


class PPTXReportBuilder:
    """Clase principal para orquestar la creación del reporte."""
    
    def __init__(self, survey, analysis_data, kpi_avg, **kwargs):
        self.survey = survey
        self.analysis_data = analysis_data
        self.kpi_avg = kpi_avg
        self.kwargs = kwargs
        self.prs = make_presentation()
        # Importante: python-pptx crea presentaciones con tamaño default (10"x7.5").
        # Nuestro layout usa widescreen (13.33"x7.5"); si no fijamos dimensiones,
        # los elementos se salen del slide al renderizar/abrir el PPTX.
        self.prs.slide_width = PPTXTheme.SLIDE_WIDTH
        self.prs.slide_height = PPTXTheme.SLIDE_HEIGHT
        self.slide_width, self.slide_height = _safe_slide_dims(self.prs)
        self.generated_at = timezone.now().strftime('%Y-%m-%d')
        self.company_name = getattr(settings, 'COMPANY_NAME', 'Byteneko SaaS')

    def build(self) -> io.BytesIO:
        """Construye todo el reporte y retorna el buffer."""
        self._create_cover_slide()
        
        if self.kwargs.get('include_kpis', True):
            self._create_executive_summary()
            
        # Opcional: Tabla resumen (podríamos integrarla o mantenerla separada, mantengamos simple por ahora)
        if self.kwargs.get('include_table', True):
           self._create_summary_table_slide()

        for item in self.analysis_data:
            self._create_detail_slide(item)

        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output

    def _create_cover_slide(self):
        """Portada estilo 'Hero' limpia."""
        # Layout vacío
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        composer = SlideComposer(slide, self.slide_width, self.slide_height)
        
        # Pinta medio fondo (diagonal o bloque lateral) para dar estilo
        # Bloque azul a la izquierda
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.5), self.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = PPTXTheme.BRAND_BLUE
        bg_shape.line.fill.background()
        
        # Logo placeholder (texto por ahora)
        tb_logo = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3.5), Inches(0.5))
        _apply_text_fitting(tb_logo.text_frame, margin=Inches(0.03))
        p_logo = tb_logo.text_frame.paragraphs[0]
        p_logo.text = self.company_name.upper()
        p_logo.font.color.rgb = PPTXTheme.WHITE
        p_logo.font.bold = True
        p_logo.font.size = Pt(14)
        p_logo.font.name = PPTXTheme.FONT_FAMILY
        
        # Título del Reporte (en la zona blanca para contraste limpio)
        title_text = getattr(self.survey, 'title', None) or (self.survey.get('title') if isinstance(self.survey, dict) else 'Reporte de Resultados')
        
        # Ajuste para evitar overflow en títulos largos
        left_content = Inches(5.0)
        top_content = Inches(2.5)
        
        # Reducir ancho para forzar margen derecho (13.33 - 5.0 - 1.0 margin = 7.33 max)
        # Usamos 7.0 para seguridad
        tb_title = slide.shapes.add_textbox(left_content, top_content, Inches(7.0), Inches(2.0))
        tf_title = tb_title.text_frame
        _apply_text_fitting(tf_title)
        
        p_t = tf_title.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = PPTXTheme.FONT_FAMILY
        p_t.font.bold = True
        p_t.font.color.rgb = PPTXTheme.DARK_INK
        p_t.font.size = Pt(44) # Max size
        
        # Subtítulo / Metadatos
        user_name = self._get_user_name()
        period = self._get_period_text()
        
        meta_text = f"Generado por: {user_name}\nFecha: {self.generated_at}\nPeríodo: {period}"
        tb_meta = slide.shapes.add_textbox(left_content, top_content + Inches(2.5), Inches(7.0), Inches(1.5))
        _apply_text_fitting(tb_meta.text_frame)
        p_meta = tb_meta.text_frame.paragraphs[0]
        p_meta.text = meta_text
        p_meta.font.size = Pt(14)
        p_meta.font.color.rgb = PPTXTheme.MUTED_INK
        p_meta.font.name = PPTXTheme.FONT_FAMILY
        
        # Decoración visual en la parte azul
        # Un círculo translúcido
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.0), Inches(5.5), Inches(3.0), Inches(3.0))
        oval.fill.solid()
        oval.fill.fore_color.rgb = PPTXTheme.WHITE
        oval.fill.transparency = 0.9
        oval.line.fill.background()

    def _create_executive_summary(self):
        """Resumen ejecutivo con KPIs principales."""
        layout = self.prs.slide_layouts[6] 
        slide = self.prs.slides.add_slide(layout)
        composer = SlideComposer(slide, self.slide_width, self.slide_height)
        composer.add_background_shape()
        
        composer.add_header("Resumen Ejecutivo", "Principales métricas y hallazgos")
        
        # KPI Cards Row
        start_y = Inches(1.8)
        gap = Inches(0.4)
        card_w = Inches(3.5)
        card_h = Inches(2.0)
        
        # 1. Total Respuestas
        total = self.kwargs.get('total_responses', 0)
        composer.add_card(Inches(0.8), start_y, card_w, card_h)
        composer.add_metric_big(Inches(1.0), start_y + Inches(0.2), "Total Respuestas", str(total), width=card_w - Inches(0.4))
        
        # 2. Satisfacción / Score General
        composer.add_card(Inches(0.8) + card_w + gap, start_y, card_w, card_h)
        composer.add_metric_big(Inches(1.0) + card_w + gap, start_y + Inches(0.2), "Score General", f"{self.kpi_avg:.1f}/10", width=card_w - Inches(0.4))
        
        # 3. NPS (si existe)
        nps_data = self.kwargs.get('nps_data') or {}
        nps_score = nps_data.get('score') if isinstance(nps_data, dict) else self.kwargs.get('nps_score')
        
        if nps_score is not None:
            composer.add_card(Inches(0.8) + (card_w + gap) * 2, start_y, card_w, card_h)
            composer.add_metric_big(Inches(1.0) + (card_w + gap) * 2, start_y + Inches(0.2), "NPS", str(nps_score), width=card_w - Inches(0.4))
            
        # Top Hallazgos (Texto)
        # Buscar insights con mood critico o excelente
        top_insights = [
            i for i in self.analysis_data 
            if i.get('insight_data', {}).get('mood') in ['CRITICO', 'EXCELENTE']
        ]
        
        list_y = start_y + card_h + Inches(0.5)
        
        if top_insights:
            tb_list = slide.shapes.add_textbox(Inches(0.8), list_y, self.slide_width - Inches(1.6), Inches(3.0))
            tf = tb_list.text_frame
            _apply_text_fitting(tf)
            p_head = tf.paragraphs[0]
            p_head.text = "Puntos de Atención:"
            p_head.font.bold = True
            p_head.font.size = Pt(14)
            p_head.font.color.rgb = PPTXTheme.DARK_INK
            
            for item in top_insights[:4]:
                insight = item.get('insight_data', {})
                mood = insight.get('mood', 'NEUTRO')
                txt = item.get('text', '')
                if isinstance(txt, str) and len(txt) > 220:
                    txt = txt[:217] + "..."
                p_i = tf.add_paragraph()
                p_i.text = f"• [{mood}] {txt}"
                p_i.font.size = Pt(12)
                p_i.space_before = Pt(6)

    def _create_summary_table_slide(self):
        """Slide con tabla de resumen de preguntas."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        composer = SlideComposer(slide, self.slide_width, self.slide_height)
        composer.add_header("Detalle de Resultados", "Vista rápida por pregunta")
        
        # Implementación simplificada de tabla
        # Calcular cuántas filas caben sin salirse de la diapositiva
        available_height = self.slide_height - Inches(2.0) # top 1.5 + margen inferior 0.5
        max_rows_fit = int(available_height / Inches(0.4)) - 1  # -1 para header
        rows = min(len(self.analysis_data), max_rows_fit, 12)
        if rows == 0: return

        cols = 3
        # table dimensions
        left = Inches(0.8)
        top = Inches(1.5)
        width = self.slide_width - Inches(1.6)
        height = Inches(0.4 * (rows + 1))

        
        table_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
        
        # Headers
        headers = ["Pregunta", "Métrica Clave", "Score/Valor"]
        # Set column widths roughly
        table_shape.columns[0].width = Inches(6.5)
        table_shape.columns[1].width = Inches(3.0)
        table_shape.columns[2].width = Inches(2.0)
        
        for idx, h_text in enumerate(headers):
            cell = table_shape.cell(0, idx)
            cell.text = h_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = PPTXTheme.mix_color(PPTXTheme.BRAND_BLUE, 0.1)
            _apply_text_fitting(cell.text_frame, margin=Inches(0.03))
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = PPTXTheme.BRAND_BLUE

        # Data
        for i, item in enumerate(self.analysis_data[:rows]):
            row_idx = i + 1
            # Col 1: Texto
            item_text = item.get('text', '')
            if len(item_text) > 80: item_text = item_text[:77] + "..."
            cell_q = table_shape.cell(row_idx, 0)
            cell_q.text = f"{item.get('order', '')}. {item_text}"
            _apply_text_fitting(cell_q.text_frame, margin=Inches(0.03))
            cell_q.text_frame.paragraphs[0].font.size = Pt(10)
            
            # Col 2 & 3: Métrica
            metric_label, metric_val = self._get_metric_summary(item)
            
            cell_m = table_shape.cell(row_idx, 1)
            cell_m.text = metric_label
            _apply_text_fitting(cell_m.text_frame, margin=Inches(0.03))
            cell_m.text_frame.paragraphs[0].font.size = Pt(10)
            
            cell_v = table_shape.cell(row_idx, 2)
            cell_v.text = str(metric_val)
            _apply_text_fitting(cell_v.text_frame, margin=Inches(0.03))
            cell_v.text_frame.paragraphs[0].font.size = Pt(10)
            cell_v.text_frame.paragraphs[0].font.bold = True

    def _create_detail_slide(self, item: dict):
        """Slide de detalle para una pregunta individual."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        composer = SlideComposer(slide, self.slide_width, self.slide_height)
        
        # Título corto
        title_full = f"{item.get('order', '')}. {item.get('text', '')}"
        composer.add_header(title_full if len(title_full) < 90 else title_full[:87] + "...")
        
        # Layout:
        # Columna Izquierda: Métricas Clave (30% ancho)
        # Columna Derecha: Gráfico (70% ancho)
        # Abajo: Narrativa
        
        insight = item.get('insight_data') or {}
        q_type = item.get('type')
        is_text_like = _is_text_like_question(item)
        
        # --- Izquierda ---
        col_l_x = Inches(0.5)
        col_l_y = Inches(1.5)
        col_l_w = Inches(3.5)
        
        # Tarjeta de métrica principal
        composer.add_card(col_l_x, col_l_y, col_l_w, Inches(2.0))

        if is_text_like:
            total_txt = item.get('total_responses', item.get('total_respuestas'))
            total_txt = int(total_txt) if isinstance(total_txt, (int, float)) else (len(item.get('top_responses') or item.get('samples_texto') or []) or 0)
            composer.add_metric_big(
                col_l_x + Inches(0.2),
                col_l_y + Inches(0.2),
                "Total Comentarios",
                str(total_txt),
                None,
                width=col_l_w - Inches(0.4),
            )

            topics = []
            if isinstance(insight, dict):
                topics = insight.get('topics') or []
            if topics:
                tb_topics = slide.shapes.add_textbox(col_l_x + Inches(0.2), col_l_y + Inches(1.25), col_l_w - Inches(0.4), Inches(0.7))
                _apply_text_fitting(tb_topics.text_frame)
                tf = tb_topics.text_frame
                p0 = tf.paragraphs[0]
                p0.text = "Temas principales"
                p0.font.name = PPTXTheme.FONT_FAMILY
                p0.font.size = Pt(10)
                p0.font.bold = True
                p0.font.color.rgb = PPTXTheme.MUTED_INK
                p1 = tf.add_paragraph()
                p1.text = ", ".join([_soft_wrap_text(str(t)) for t in topics[:3]])
                p1.font.name = PPTXTheme.FONT_FAMILY
                p1.font.size = Pt(11)
                p1.font.color.rgb = PPTXTheme.DARK_INK
        else:
            # Contenido métrica (numérica/categórica)
            lbl, val = self._get_metric_summary(item)
            composer.add_metric_big(
                col_l_x + Inches(0.2),
                col_l_y + Inches(0.2),
                lbl,
                str(val),
                insight.get('trend_delta') if isinstance(insight, dict) else None,
                width=col_l_w - Inches(0.4),
            )
        
        # --- Derecha (Gráfico) ---
        col_r_x = Inches(4.5)
        col_r_y = Inches(1.5)
        # Asegurar que no se salga del margen derecho: slide_width=13.33, necesitamos margen de 0.5
        col_r_w = self.slide_width - col_r_x - Inches(0.5)
        col_r_h = Inches(3.5)
        
        # Contenido derecha: gráfico o lista de comentarios
        composer.add_card(col_r_x, col_r_y, col_r_w, col_r_h, bg_color=PPTXTheme.WHITE)

        if is_text_like:
            # Encabezado interno
            tb_head = slide.shapes.add_textbox(col_r_x + Inches(0.25), col_r_y + Inches(0.2), col_r_w - Inches(0.5), Inches(0.35))
            _apply_text_fitting(tb_head.text_frame, margin=Inches(0.03))
            p_h = tb_head.text_frame.paragraphs[0]
            p_h.text = "Comentarios destacados"
            p_h.font.name = PPTXTheme.FONT_FAMILY
            p_h.font.size = Pt(12)
            p_h.font.bold = True
            p_h.font.color.rgb = PPTXTheme.DARK_INK

            # Lista de comentarios (muestras)
            samples = item.get('top_responses') or item.get('samples_texto') or item.get('samples') or []
            samples = [s for s in samples if isinstance(s, str) and s.strip()]
            max_items = 7
            tb_list = slide.shapes.add_textbox(col_r_x + Inches(0.25), col_r_y + Inches(0.6), col_r_w - Inches(0.5), col_r_h - Inches(0.85))
            tf_list = tb_list.text_frame
            _apply_text_fitting(tf_list)
            tf_list.clear()

            if not samples:
                p0 = tf_list.paragraphs[0]
                p0.text = "Sin comentarios para mostrar"
                p0.font.name = PPTXTheme.FONT_FAMILY
                p0.font.size = Pt(12)
                p0.font.color.rgb = PPTXTheme.MUTED_INK
            else:
                for idx, raw in enumerate(samples[:max_items]):
                    text = " ".join(raw.split())
                    if len(text) > 170:
                        text = text[:167] + "..."
                    text = _soft_wrap_text(text)
                    p = tf_list.paragraphs[0] if idx == 0 else tf_list.add_paragraph()
                    p.text = f"• {text}"
                    p.font.name = PPTXTheme.FONT_FAMILY
                    p.font.size = Pt(12)
                    p.font.color.rgb = PPTXTheme.DARK_INK
                    p.space_before = Pt(4)
        else:
            # Generar e insertar gráfico
            if self.kwargs.get('include_charts', True):
                self._insert_chart(slide, item, col_r_x, col_r_y, col_r_w, col_r_h)

        # --- Abajo (Narrativa) ---
        narrative = insight.get('narrative')
        if narrative:
            # Asegurar que la caja de narrativa respete margen inferior
            # slide_height=7.5, queremos terminar en 7.0 max
            narr_y = Inches(5.2)
            narr_h = min(Inches(1.6), self.slide_height - narr_y - Inches(0.3))
            composer.add_narrative_box(
                col_l_x, 
                narr_y, 
                self.slide_width - Inches(1.0), 
                narr_h, 
                narrative
            )

    def _insert_chart(self, slide, item, x, y, w, h):
        """Genera e inserta la imagen del gráfico."""
        try:
            from core.utils.charts import ChartGenerator
            labels = item.get('chart_labels') or []
            counts = item.get('chart_data') or []
            
            if not labels or not counts:
                return

            chart_b64 = None
            q_type = item.get('type')
            
            # Lógica simple de selección de chart
            if q_type in ['single', 'multi', 'radio', 'select', 'boolean']:
                if (item.get('tipo_display') == 'doughnut') or (len(labels) <= 5):
                     chart_b64 = ChartGenerator.generate_pie_chart(labels, counts, title='', dark_mode=False)
                else:
                     chart_b64 = ChartGenerator.generate_horizontal_bar_chart(labels, counts, title='', dark_mode=False)
            elif q_type in ['scale', 'number', 'numeric', 'rating']:
                 chart_b64 = ChartGenerator.generate_horizontal_bar_chart(labels, counts, title='', dark_mode=False)
            
            if chart_b64:
                img_bytes = base64.b64decode(chart_b64)
                # Insertar sin deformar (mantener proporción) dentro del contenedor.
                _add_picture_contain(slide, img_bytes, x, y, w, h, padding=Inches(0.12))

        except Exception:
            logger.exception("Error generando chart en PPTX")

    def _get_metric_summary(self, item) -> Tuple[str, Any]:
        """Devuelve (Label, Valor) representativo para la pregunta."""
        insight = item.get('insight_data') or {}
        q_type = item.get('type')
        
        if q_type in ['scale', 'number', 'numeric', 'rating']:
            avg = insight.get('avg') if insight.get('avg') is not None else insight.get('average')
            val = f"{float(avg):.1f}" if avg is not None else "N/A"
            return ("Promedio", val)
        
        if q_type in ['single', 'multi', 'radio', 'select', 'boolean']:
            top = insight.get('top_option')
            if top and top.get('option'):
                # Truncar opción si es muy larga
                opt = str(top['option'])
                if len(opt) > 20: opt = opt[:17] + "..."
                pct = ""
                if insight.get('total'):
                     p = (top['count'] / insight['total']) * 100
                     pct = f" ({p:.0f}%)"
                return ("Top Opción", f"{opt}{pct}")
            return ("Top Opción", "N/A")

        if q_type in ['text', 'comment', 'comments', 'textarea'] or (item.get('tipo_display') == 'text'):
            total = item.get('total_responses', item.get('total_respuestas'))
            if total is None:
                samples = item.get('top_responses') or item.get('samples_texto') or []
                total = len(samples) if isinstance(samples, list) else 0
            return ("Comentarios", total)
            
        return ("Respuesta", "-")

    def _get_user_name(self) -> str:
        # Priorizar argumento explícito
        if self.kwargs.get('user_name'):
            return self.kwargs['user_name']
            
        req = self.kwargs.get('request')
        if req and hasattr(req, 'user') and req.user.is_authenticated:
            return req.user.get_full_name() or req.user.username
        
        # Fallback si no hay request pero hay pass de contexto manual
        return "Usuario del Sistema"

    def _get_period_text(self) -> str:
        s = self.kwargs.get('start_date')
        e = self.kwargs.get('end_date')
        if s and e: return f"{s} - {e}"
        if s: return f"Desde {s}"
        days = self.kwargs.get('window_days')
        if days and days != 'all': return f"Últimos {days} días"
        return "Todo el histórico"


# --- API Pública de Compatibilidad ---

def generate_full_pptx_report(survey, analysis_data, kpi_satisfaction_avg: float = 0.0, **kwargs) -> io.BytesIO:
    """Función de entrada principal para generar el reporte PPTX."""
    builder = PPTXReportBuilder(survey, analysis_data, kpi_satisfaction_avg, **kwargs)
    return builder.build()


class PPTXReportGenerator:
    """Clase wrapper legacy para compatibilidad con código existente."""
    @classmethod
    def generate(cls, survey, analysis_data, kpi_satisfaction_avg=0, **kwargs):
        return generate_full_pptx_report(survey, analysis_data, kpi_satisfaction_avg, **kwargs)